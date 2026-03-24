#!/usr/bin/env python3
"""
NotebookLM Slide Deck Converter
Converts PDF slide decks from NotebookLM into PPTX, PNG, or HTML formats.

Usage:
    python convert.py --input slides.pdf --output slides.pptx --format pptx
    python convert.py --input slides.pdf --output ./slides/ --format png
    python convert.py --input slides.pdf --output slides.html --format html
    python convert.py --input slides.pdf --output slides.pptx --format pptx --dpi 300
"""

import argparse
import os
import sys
import tempfile
from pathlib import Path


def convert_to_pptx(pdf_path: str, output_path: str, dpi: int = 200):
    """Convert PDF to PPTX with full-bleed slide images."""
    from pdf2image import convert_from_path
    from pptx import Presentation
    from pptx.util import Inches

    print(f"Converting PDF to images at {dpi} DPI...")
    images = convert_from_path(pdf_path, dpi=dpi)
    print(f"Got {len(images)} slides")

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    for i, img in enumerate(images):
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
        tmp = tempfile.NamedTemporaryFile(suffix=".png", delete=False)
        img.save(tmp.name, "PNG")
        slide.shapes.add_picture(
            tmp.name, Inches(0), Inches(0), prs.slide_width, prs.slide_height
        )
        os.unlink(tmp.name)
        print(f"  Slide {i + 1}/{len(images)} done")

    prs.save(output_path)
    size_mb = os.path.getsize(output_path) / 1024 / 1024
    print(f"\n✅ Saved: {output_path} ({size_mb:.1f} MB, {len(images)} slides)")


def convert_to_png(pdf_path: str, output_dir: str, dpi: int = 300):
    """Convert PDF to individual PNG slide images."""
    from pdf2image import convert_from_path

    os.makedirs(output_dir, exist_ok=True)

    print(f"Converting PDF to images at {dpi} DPI...")
    images = convert_from_path(pdf_path, dpi=dpi)
    print(f"Got {len(images)} slides")

    paths = []
    for i, img in enumerate(images):
        filename = f"slide_{i + 1:02d}.png"
        filepath = os.path.join(output_dir, filename)
        img.save(filepath, "PNG")
        size_kb = os.path.getsize(filepath) / 1024
        print(f"  {filename} ({size_kb:.0f} KB)")
        paths.append(filepath)

    print(f"\n✅ Exported {len(images)} slides to: {output_dir}")
    return paths


def convert_to_html(pdf_path: str, output_path: str, dpi: int = 200):
    """Convert PDF to a self-contained HTML slideshow."""
    import base64
    from pdf2image import convert_from_path

    print(f"Converting PDF to images at {dpi} DPI...")
    images = convert_from_path(pdf_path, dpi=dpi)
    print(f"Got {len(images)} slides")

    # Encode images as base64
    encoded_images = []
    for i, img in enumerate(images):
        tmp = tempfile.NamedTemporaryFile(suffix=".png", delete=False)
        img.save(tmp.name, "PNG")
        with open(tmp.name, "rb") as f:
            b64 = base64.b64encode(f.read()).decode()
        encoded_images.append(b64)
        os.unlink(tmp.name)
        print(f"  Slide {i + 1}/{len(images)} encoded")

    # Build HTML slideshow
    slides_html = ""
    for i, b64 in enumerate(encoded_images):
        display = "block" if i == 0 else "none"
        slides_html += f'<div class="slide" style="display:{display}"><img src="data:image/png;base64,{b64}" alt="Slide {i+1}"></div>\n'

    html = f"""<!DOCTYPE html>
<html lang="en">
<head>
<meta charset="UTF-8">
<meta name="viewport" content="width=device-width, initial-scale=1.0">
<title>NotebookLM Slides</title>
<style>
  * {{ margin: 0; padding: 0; box-sizing: border-box; }}
  body {{ background: #1a1a2e; display: flex; flex-direction: column; align-items: center; justify-content: center; min-height: 100vh; font-family: -apple-system, sans-serif; }}
  .slide img {{ max-width: 90vw; max-height: 80vh; border-radius: 12px; box-shadow: 0 20px 60px rgba(0,0,0,0.5); }}
  .controls {{ margin-top: 24px; display: flex; gap: 16px; align-items: center; }}
  .controls button {{ padding: 12px 28px; border: none; border-radius: 8px; background: #e94560; color: white; font-size: 16px; cursor: pointer; transition: all 0.2s; }}
  .controls button:hover {{ background: #c23152; transform: translateY(-2px); }}
  .controls button:disabled {{ opacity: 0.3; cursor: not-allowed; transform: none; }}
  .counter {{ color: #888; font-size: 14px; min-width: 80px; text-align: center; }}
</style>
</head>
<body>
{slides_html}
<div class="controls">
  <button id="prev" onclick="navigate(-1)">← Prev</button>
  <span class="counter" id="counter">1 / {len(encoded_images)}</span>
  <button id="next" onclick="navigate(1)">Next →</button>
</div>
<script>
  let current = 0;
  const slides = document.querySelectorAll('.slide');
  const total = slides.length;
  function navigate(dir) {{
    slides[current].style.display = 'none';
    current = Math.max(0, Math.min(total - 1, current + dir));
    slides[current].style.display = 'block';
    document.getElementById('counter').textContent = (current + 1) + ' / ' + total;
    document.getElementById('prev').disabled = current === 0;
    document.getElementById('next').disabled = current === total - 1;
  }}
  document.addEventListener('keydown', e => {{
    if (e.key === 'ArrowLeft') navigate(-1);
    if (e.key === 'ArrowRight') navigate(1);
  }});
  navigate(0);
</script>
</body>
</html>"""

    with open(output_path, "w") as f:
        f.write(html)

    size_mb = os.path.getsize(output_path) / 1024 / 1024
    print(f"\n✅ Saved: {output_path} ({size_mb:.1f} MB, {len(encoded_images)} slides)")


def main():
    parser = argparse.ArgumentParser(description="Convert NotebookLM PDF slides to other formats")
    parser.add_argument("--input", "-i", required=True, help="Input PDF file path")
    parser.add_argument("--output", "-o", required=True, help="Output file/directory path")
    parser.add_argument("--format", "-f", choices=["pptx", "png", "html"], default="pptx", help="Output format (default: pptx)")
    parser.add_argument("--dpi", "-d", type=int, default=200, help="Image resolution in DPI (default: 200)")

    args = parser.parse_args()

    if not os.path.exists(args.input):
        print(f"Error: Input file not found: {args.input}")
        sys.exit(1)

    if args.format == "pptx":
        convert_to_pptx(args.input, args.output, args.dpi)
    elif args.format == "png":
        convert_to_png(args.input, args.output, args.dpi)
    elif args.format == "html":
        convert_to_html(args.input, args.output, args.dpi)


if __name__ == "__main__":
    main()
